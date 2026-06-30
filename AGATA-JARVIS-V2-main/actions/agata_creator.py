import json
import os
import re
import subprocess
import sys
import time
import uuid
from pathlib import Path

import requests

OPENCODE_API_KEY = "sk-6uTKF9DpKWQ8yxQoVi8b6hMSvfwjjBtibv786RgzWohnad0z9sdGV9fnKjTwJTNS"
OPENCODE_BASE_URL = "https://opencode.ai/zen/v1"
OPENCODE_MODEL = "nemotron-3-super-free"
OPENCODE_FALLBACK_MODEL = "deepseek-v4-flash-free"

DESKTOP = Path.home() / "Desktop"
AGATA_FOLDER = Path(r"C:\Users\ASUS\OneDrive\Documenti")


def _ensure_folder():
    AGATA_FOLDER.mkdir(parents=True, exist_ok=True)


def _call_opencode(prompt: str, max_tokens: int = 4000) -> str:
    last_error = ""
    for model in [OPENCODE_MODEL, OPENCODE_FALLBACK_MODEL]:
        headers = {
            "Authorization": f"Bearer {OPENCODE_API_KEY}",
            "Content-Type": "application/json",
        }
        payload = {
            "model": model,
            "messages": [
                {
                    "role": "system",
                    "content": (
                        "You are an elite document designer and content creator. "
                        "You craft beautifully structured, professional content for documents and presentations. "
                        "For Word documents: include rich formatting suggestions, elegant structure, professional tone. "
                        "For PowerPoint: break content into slide-by-slide format with design notes. "
                        "Use professional language with a creative, elegant flair. "
                        "Respond in the same language as the prompt. "
                        "IMPORTANT: Output ONLY the final content. Do NOT include your thinking process, analysis, "
                        "or reasoning. Start directly with the document/presentation content."
                    ),
                },
                {"role": "user", "content": prompt},
            ],
            "max_tokens": max_tokens,
            "temperature": 0.7,
        }
        try:
            r = requests.post(
                f"{OPENCODE_BASE_URL}/chat/completions",
                headers=headers,
                json=payload,
                timeout=60,
            )
            if r.status_code == 200:
                msg = r.json()["choices"][0]["message"]
                content = msg.get("content", "") or msg.get("reasoning_content", "")
                if content and len(content.strip()) > 20:
                    return content
                last_error = f"[OpenCode] Model {model} returned empty content"
            else:
                last_error = f"[OpenCode Error {r.status_code}]: {r.text[:300]}"
        except Exception as e:
            last_error = f"[OpenCode Connection Error]: {str(e)[:200]}"
    return last_error or "[OpenCode Error]: All models failed"


def _call_ollama(prompt: str, max_tokens: int = 4000) -> str:
    try:
        r = requests.post(
            "http://localhost:11434/api/generate",
            json={
                "model": "phi:2.7b",
                "prompt": prompt,
                "stream": False,
                "options": {
                    "num_predict": max_tokens,
                    "temperature": 0.7,
                },
            },
            timeout=180,
        )
        if r.status_code == 200:
            data = r.json()
            content = data.get("response", "") or data.get("thinking", "") or data.get("message", {}).get("content", "")
            return content
        return f"[Ollama Error {r.status_code}]: {r.text[:300]}"
    except Exception as e:
        return f"[Ollama Connection Error]: {str(e)[:200]}"


def _call_gemini_fast(prompt: str, api_key: str) -> str:
    from core.config import gemini_generate
    from core.ai_providers import generate_gemini_cli
    try:
        cli_result = generate_gemini_cli(prompt, max_tokens=2000)
        if cli_result and not cli_result.startswith("[Gemini CLI"):
            return cli_result
    except Exception:
        pass
    try:
        return gemini_generate(prompt, model="gemini-2.0-flash", max_tokens=2000)
    except Exception as e:
        return f"[Gemini Fast Error]: {str(e)[:200]}"


# ─── COLOR PALETTES ────────────────────────────────────────────────────

PALETTES = {
    "elegant": {
        "primary": "#2C3E50", "secondary": "#E74C3C", "accent": "#3498DB",
        "light": "#ECF0F1", "dark": "#1A252F", "text": "#2C3E50",
    },
    "pastel": {
        "primary": "#6C5CE7", "secondary": "#FD79A8", "accent": "#00CEC9",
        "light": "#F8F5FF", "dark": "#2D1B69", "text": "#4A3F6B",
    },
    "corporativo": {
        "primary": "#1B4F72", "secondary": "#2874A6", "accent": "#D4AC0D",
        "light": "#EBF5FB", "dark": "#0B2D47", "text": "#1B4F72",
    },
    "moderno": {
        "primary": "#0A0A0A", "secondary": "#FF6B6B", "accent": "#4ECDC4",
        "light": "#F7F7F7", "dark": "#0A0A0A", "text": "#2D2D2D",
    },
    "rosa": {
        "primary": "#E91E63", "secondary": "#9C27B0", "accent": "#FFD54F",
        "light": "#FCE4EC", "dark": "#880E4F", "text": "#4A0072",
    },
    "nordico": {
        "primary": "#5E81AC", "secondary": "#BF616A", "accent": "#A3BE8C",
        "light": "#ECEFF4", "dark": "#2E3440", "text": "#3B4252",
    },
}


def _design_word(content: str, title: str, palette: str, player=None, speak=None) -> str:
    try:
        from docx import Document
        from docx.shared import Inches, Pt, RGBColor, Cm, Emu
        from docx.enum.text import WD_ALIGN_PARAGRAPH, WD_LINE_SPACING
        from docx.enum.table import WD_TABLE_ALIGNMENT
        from docx.enum.section import WD_ORIENT
        from docx.oxml.ns import qn, nsdecls
        from docx.oxml import parse_xml
    except ImportError:
        raise ImportError("python-docx no esta instalado. Ejecuta: pip install python-docx")

    from .utils_diseno import buscar_y_descargar_imagen

    _ensure_folder()
    pal = PALETTES.get(palette, PALETTES["elegant"])
    primary_rgb = _hex_to_rgb(pal["primary"])
    secondary_rgb = _hex_to_rgb(pal["secondary"])
    accent_rgb = _hex_to_rgb(pal["accent"])
    text_rgb = _hex_to_rgb(pal["text"])
    light_rgb = _hex_to_rgb(pal["light"])
    dark_rgb = _hex_to_rgb(pal["dark"])

    filename = f"Agata_{title.replace(' ', '_')[:40]}_{uuid.uuid4().hex[:6]}.docx"
    filepath = AGATA_FOLDER / filename

    doc = Document()

    # configure default style
    style = doc.styles["Normal"]
    font = style.font
    font.name = "Calibri"
    font.size = Pt(11)
    font.color.rgb = RGBColor(*text_rgb)
    pf = style.paragraph_format
    pf.space_after = Pt(8)
    pf.line_spacing = 1.3

    # configure heading styles
    for i, (h_style, h_size, h_color) in enumerate([
        ("Heading 1", 22, primary_rgb),
        ("Heading 2", 16, secondary_rgb),
        ("Heading 3", 13, accent_rgb),
    ]):
        try:
            hs = doc.styles[h_style]
            hs.font.size = Pt(h_size)
            hs.font.bold = True
            hs.font.color.rgb = RGBColor(*h_color)
            hs.font.name = "Calibri Light"
            hs.paragraph_format.space_before = Pt(18 if i < 2 else 12)
            hs.paragraph_format.space_after = Pt(8)
        except Exception:
            pass

    for section in doc.sections:
        section.top_margin = Cm(2.5)
        section.bottom_margin = Cm(2.5)
        section.left_margin = Cm(2.8)
        section.right_margin = Cm(2.8)

    # ── COVER PAGE ──
    for _ in range(5):
        doc.add_paragraph("")

    # decorative top line
    line_p = doc.add_paragraph()
    line_p.alignment = WD_ALIGN_PARAGRAPH.CENTER
    lr = line_p.add_run("━" * 40)
    lr.font.color.rgb = RGBColor(*accent_rgb)
    lr.font.size = Pt(10)

    title_p = doc.add_paragraph()
    title_p.alignment = WD_ALIGN_PARAGRAPH.CENTER
    tr = title_p.add_run(title.upper())
    tr.font.size = Pt(34)
    tr.font.bold = True
    tr.font.color.rgb = RGBColor(*primary_rgb)
    tr.font.name = "Calibri Light"

    subtitle_p = doc.add_paragraph()
    subtitle_p.alignment = WD_ALIGN_PARAGRAPH.CENTER
    sr = subtitle_p.add_run("Disenado por Agata  |  JARVIS Industries")
    sr.font.size = Pt(12)
    sr.font.color.rgb = RGBColor(*secondary_rgb)
    sr.font.italic = True

    date_p = doc.add_paragraph()
    date_p.alignment = WD_ALIGN_PARAGRAPH.CENTER
    dr = date_p.add_run(time.strftime("%d de %B de %Y"))
    dr.font.size = Pt(10)
    dr.font.color.rgb = RGBColor(150, 150, 150)

    # bottom decorative line
    line_p2 = doc.add_paragraph()
    line_p2.alignment = WD_ALIGN_PARAGRAPH.CENTER
    lr2 = line_p2.add_run("━" * 40)
    lr2.font.color.rgb = RGBColor(*accent_rgb)
    lr2.font.size = Pt(10)

    doc.add_page_break()

    # ── CONTENT WITH TABLES / IMAGES ──
    lines = content.strip().split("\n")
    i = 0
    while i < len(lines):
        line = lines[i].strip()
        if not line:
            i += 1
            continue

        # heading detection
        if line.startswith("# ") or line.startswith("## ") or line.startswith("### "):
            level = len(line) - len(line.lstrip("#"))
            heading_text = re.sub(r"^#+\s*", "", line)
            p = doc.add_paragraph()
            p.style = doc.styles[f"Heading {min(level, 3)}"]
            run = p.add_run(heading_text)

        # image insertion
        elif line.startswith("[IMAGE:") or line.startswith("[IMAGEN:"):
            kw_match = re.search(r"\[(?:IMAGE|IMAGEN):\s*(.+?)\]", line, re.IGNORECASE)
            if kw_match and speak:
                pass
            if kw_match:
                keyword = kw_match.group(1).strip()
                player.write_log(f"[Agata] Buscando imagen: {keyword}")
                img_path = buscar_y_descargar_imagen(keyword)
                if img_path:
                    try:
                        p = doc.add_paragraph()
                        p.alignment = WD_ALIGN_PARAGRAPH.CENTER
                        run = p.add_run()
                        run.add_picture(img_path, width=Inches(4.5))
                        doc.add_paragraph("")
                    except Exception:
                        pass

        # table detection
        elif line.startswith("[TABLE]") or line.startswith("[TABLA]"):
            table_lines = []
            i += 1
            while i < len(lines) and not lines[i].strip().startswith("[ENDTABLE]") and not lines[i].strip().startswith("[FIN_TABLA]"):
                table_lines.append(lines[i].strip())
                i += 1
            if table_lines:
                _build_table(doc, table_lines, pal)

        # bullet points
        elif line.startswith("- "):
            bullet_text = line[2:]
            p = doc.add_paragraph(style="List Bullet")
            p.clear()
            run = p.add_run(bullet_text)
            run.font.size = Pt(11)
            run.font.color.rgb = RGBColor(*text_rgb)

        # numbered list
        elif re.match(r"^\d+[\.\)]\s", line):
            p = doc.add_paragraph(style="List Number")
            p.clear()
            run = p.add_run(re.sub(r"^\d+[\.\)]\s*", "", line))
            run.font.size = Pt(11)
            run.font.color.rgb = RGBColor(*text_rgb)

        # blockquote
        elif line.startswith("> "):
            quote_text = line[2:]
            p = doc.add_paragraph()
            p.paragraph_format.left_indent = Cm(1.5)
            p.paragraph_format.space_before = Pt(8)
            p.paragraph_format.space_after = Pt(8)
            run = p.add_run(quote_text)
            run.font.size = Pt(11)
            run.font.italic = True
            run.font.color.rgb = RGBColor(*secondary_rgb)

        # separator / divider
        elif line.startswith("---") or line.startswith("==="):
            p = doc.add_paragraph()
            p.alignment = WD_ALIGN_PARAGRAPH.CENTER
            r = p.add_run("─" * 50)
            r.font.color.rgb = RGBColor(*accent_rgb)
            r.font.size = Pt(8)

        # regular paragraph
        else:
            p = doc.add_paragraph()
            run = p.add_run(line)
            run.font.size = Pt(11)
            run.font.color.rgb = RGBColor(*text_rgb)

        i += 1

    # ── FOOTER WITH PAGE NUMBERS ──
    doc.add_paragraph("")
    p = doc.add_paragraph()
    p.alignment = WD_ALIGN_PARAGRAPH.CENTER
    r = p.add_run("─" * 40)
    r.font.color.rgb = RGBColor(*accent_rgb)
    r.font.size = Pt(8)

    footer_para = doc.add_paragraph()
    footer_para.alignment = WD_ALIGN_PARAGRAPH.CENTER
    fr = footer_para.add_run("Creado con ❤ por Agata  |  JARVIS Industries")
    fr.font.size = Pt(8)
    fr.font.color.rgb = RGBColor(*secondary_rgb)
    fr.font.italic = True

    # add page numbers to footer
    for section in doc.sections:
        footer = section.footer
        footer.is_linked_to_previous = False
        fp = footer.paragraphs[0]
        fp.alignment = WD_ALIGN_PARAGRAPH.CENTER
        fpr = fp.add_run()
        fldChar1 = parse_xml(f'<w:fldChar {nsdecls("w")} w:fldCharType="begin"/>')
        fpr._r.append(fldChar1)
        instrText = parse_xml(f'<w:instrText {nsdecls("w")} xml:space="preserve"> PAGE </w:instrText>')
        fpr._r.append(instrText)
        fldChar2 = parse_xml(f'<w:fldChar {nsdecls("w")} w:fldCharType="end"/>')
        fpr._r.append(fldChar2)
        fpr.font.size = Pt(8)
        fpr.font.color.rgb = RGBColor(150, 150, 150)

    doc.save(str(filepath))
    return str(filepath)


def _build_table(doc, lines: list[str], pal: dict):
    from docx.shared import Inches, Pt, RGBColor, Cm
    from docx.enum.text import WD_ALIGN_PARAGRAPH
    from docx.enum.table import WD_TABLE_ALIGNMENT
    from docx.oxml.ns import qn, nsdecls
    from docx.oxml import parse_xml

    primary_rgb = _hex_to_rgb(pal["primary"])
    light_rgb = _hex_to_rgb(pal["light"])
    dark_rgb = _hex_to_rgb(pal["dark"])
    text_rgb = _hex_to_rgb(pal["text"])

    rows_data = []
    for line in lines:
        if line.startswith("|") and line.endswith("|"):
            cells = [c.strip() for c in line.strip("|").split("|")]
            rows_data.append(cells)

    if not rows_data:
        return

    num_cols = max(len(r) for r in rows_data)
    table = doc.add_table(rows=len(rows_data), cols=num_cols)
    table.style = "Table Grid"
    table.alignment = WD_TABLE_ALIGNMENT.CENTER

    for row_idx, row_data in enumerate(rows_data):
        for col_idx in range(num_cols):
            cell = table.cell(row_idx, col_idx)
            cell_text = row_data[col_idx] if col_idx < len(row_data) else ""
            cell.paragraphs[0].clear()
            run = cell.paragraphs[0].add_run(cell_text)

            if row_idx == 0:
                run.font.bold = True
                run.font.size = Pt(10)
                run.font.color.rgb = RGBColor(255, 255, 255)
                shading = parse_xml(
                    f'<w:shd {nsdecls("w")} w:fill="{pal["primary"].lstrip("#")}" w:val="clear"/>'
                )
                cell._tc.get_or_add_tcPr().append(shading)
            else:
                run.font.size = Pt(10)
                run.font.color.rgb = RGBColor(*text_rgb)
                if row_idx % 2 == 0:
                    shading = parse_xml(
                        f'<w:shd {nsdecls("w")} w:fill="{pal["light"].lstrip("#")}" w:val="clear"/>'
                    )
                    cell._tc.get_or_add_tcPr().append(shading)

            cell.paragraphs[0].alignment = WD_ALIGN_PARAGRAPH.LEFT

    doc.add_paragraph("")


def _parse_slides_from_content(content: str) -> list[dict]:
    import re

    slides = []
    blocks = content.split("\n---\n")
    if len(blocks) == 1:
        blocks = content.split("\n\n##")

    for block in blocks:
        block = block.strip()
        if not block:
            continue

        lines = block.split("\n")
        slide_title = ""
        keyword = ""
        slide_lines = []

        for i, line in enumerate(lines):
            stripped = line.strip()
            if not stripped:
                continue

            if i == 0 and (stripped.startswith("## ") or stripped.startswith("# ")):
                raw = re.sub(r"^#+\s*", "", stripped)
                kw_match = re.search(r"\[keyword:\s*(.+?)\]", raw, re.IGNORECASE)
                if kw_match:
                    keyword = kw_match.group(1).strip()
                    raw = re.sub(r"\s*\[keyword:.*?\]", "", raw, flags=re.IGNORECASE).strip()
                slide_title = raw
                slide_lines.append(f"## {raw}")
            else:
                slide_lines.append(stripped)

        if not slide_title:
            for line in lines:
                stripped = line.strip()
                if stripped and not stripped.startswith("-") and not stripped.startswith("**"):
                    slide_title = stripped[:80]
                    break

        if not keyword and slide_title:
            keyword = f"{slide_title} high quality background"

        slide_text = "\n".join(slide_lines) if slide_lines else block
        slides.append({"title": slide_title, "keyword": keyword, "text": slide_text})

    return slides


def _design_ppt(content: str, title: str, palette: str, player=None, speak=None) -> str:
    try:
        from pptx import Presentation
        from pptx.util import Inches, Pt
        from pptx.dml.color import RGBColor
        from pptx.enum.text import PP_ALIGN
    except ImportError:
        raise ImportError("python-pptx no esta instalado. Ejecuta: pip install python-pptx")

    from .utils_diseno import buscar_y_descargar_imagen, analizar_paleta_y_brillo

    def _overlay(slide):
        from pptx.oxml.ns import qn

        ov = slide.shapes.add_shape(1, 0, 0, prs.slide_width, prs.slide_height)
        ov.fill.solid()
        ov.fill.fore_color.rgb = RGBColor(0, 0, 0)
        ov.line.fill.background()
        sp_el = ov._element
        sf = sp_el.find('.//' + qn('a:solidFill'))
        if sf is not None:
            srgb = sf.find(qn('a:srgbClr'))
            if srgb is not None:
                alpha_attr = qn('a:alpha')
                srgb.set(alpha_attr, '50000')
        else:
            srgb = sp_el.find('.//' + qn('a:srgbClr'))
            if srgb is not None:
                srgb.set(qn('a:alpha'), '50000')
        return ov

    _ensure_folder()
    pal = PALETTES.get(palette, PALETTES["moderno"])
    filename = f"Agata_{title.replace(' ', '_')[:40]}_{uuid.uuid4().hex[:6]}.pptx"
    filepath = AGATA_FOLDER / filename

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    slides_data = _parse_slides_from_content(content)

    if not slides_data:
        slides_data = [{"title": title, "keyword": f"{title} technology abstract", "text": content}]

    # ── SLIDE 0: TITLE SLIDE ──
    if speak:
        pass
    title_keyword = f"{title} abstract technology background"
    title_image = buscar_y_descargar_imagen(title_keyword)
    title_info = analizar_paleta_y_brillo(title_image) if title_image else {
        "accent": _hex_to_rgb(pal["accent"]),
        "text_main": _hex_to_rgb(pal["light"]),
        "text_secondary": _hex_to_rgb(pal["light"]),
        "is_background_light": False,
    }

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    if title_image:
        slide.shapes.add_picture(title_image, 0, 0, prs.slide_width, prs.slide_height)

    if not title_image or title_info["is_background_light"]:
        _overlay(slide)

    txBox = slide.shapes.add_textbox(Inches(1.5), Inches(2), Inches(10), Inches(2))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = title.upper()
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = RGBColor(*title_info["text_main"])
    p.font.name = "Calibri Light"
    p.alignment = PP_ALIGN.CENTER

    sub_box = slide.shapes.add_textbox(Inches(1.5), Inches(4.2), Inches(10), Inches(1))
    sub_tf = sub_box.text_frame
    sp = sub_tf.paragraphs[0]
    sp.text = "Disenado por Agata"
    sp.font.size = Pt(18)
    sp.font.color.rgb = RGBColor(*title_info["accent"])
    sp.font.italic = True
    sp.alignment = PP_ALIGN.CENTER

    line_shape = slide.shapes.add_shape(
        1, Inches(5), Inches(3.5), Inches(3), Inches(0.03)
    )
    line_shape.fill.solid()
    line_shape.fill.fore_color.rgb = RGBColor(*title_info["accent"])
    line_shape.line.fill.background()

    # ── CONTENT SLIDES ──
    for slide_idx, sd in enumerate(slides_data):
        if speak and slide_idx % 3 == 0:
            pass

        player.write_log(f"[Agata] Buscando imagen para: {sd['keyword'][:60]}")
        img_path = buscar_y_descargar_imagen(sd["keyword"])
        info_color = analizar_paleta_y_brillo(img_path) if img_path else {
            "accent": _hex_to_rgb(pal["accent"]),
            "text_main": _hex_to_rgb(pal["dark"]),
            "text_secondary": (100, 100, 100),
            "is_background_light": True,
        }

        slide = prs.slides.add_slide(prs.slide_layouts[6])

        if img_path:
            slide.shapes.add_picture(img_path, 0, 0, prs.slide_width, prs.slide_height)

        has_text_color = info_color["text_main"]

        if not img_path or info_color["is_background_light"]:
            _overlay(slide)
            if not img_path:
                has_text_color = _hex_to_rgb(pal["light"])
                info_color["accent"] = _hex_to_rgb(pal["accent"])
                info_color["text_secondary"] = (180, 180, 180)

        lines = sd["text"].strip().split("\n")
        y_pos = 1.0 if len(lines) < 6 else 0.5

        has_title = False
        for line in lines:
            line = line.strip()
            if not line:
                y_pos += 0.3
                continue

            if (line.startswith("## ") or line.startswith("# ")) and not has_title:
                heading_text = re.sub(r"^#+\s*", "", line)
                heading_text = re.sub(r"\s*\[keyword:.*?\]", "", heading_text, flags=re.IGNORECASE).strip()
                tb = slide.shapes.add_textbox(Inches(1), Inches(y_pos), Inches(11.3), Inches(1))
                ttf = tb.text_frame
                tp = ttf.paragraphs[0]
                tp.text = heading_text
                tp.font.size = Pt(32)
                tp.font.bold = True
                tp.font.color.rgb = RGBColor(*info_color["accent"])
                tp.font.name = "Calibri Light"
                tp.alignment = PP_ALIGN.LEFT
                y_pos += 1.2
                has_title = True

            elif line.startswith("- "):
                bullet = line[2:]
                tb = slide.shapes.add_textbox(Inches(1.5), Inches(y_pos), Inches(10.3), Inches(0.5))
                ttf = tb.text_frame
                tp = ttf.paragraphs[0]
                tp.text = f"   {bullet}"
                tp.font.size = Pt(16)
                tp.font.color.rgb = RGBColor(*has_text_color)
                tp.font.name = "Calibri"
                y_pos += 0.55

            elif line.startswith("**") and line.endswith("**"):
                highlight = line.strip("*")
                tb = slide.shapes.add_textbox(Inches(1), Inches(y_pos), Inches(11.3), Inches(0.6))
                ttf = tb.text_frame
                tp = ttf.paragraphs[0]
                tp.text = highlight
                tp.font.size = Pt(20)
                tp.font.bold = True
                tp.font.color.rgb = RGBColor(*info_color["accent"])
                tp.font.name = "Calibri"
                tp.alignment = PP_ALIGN.LEFT
                y_pos += 0.8

            else:
                tb = slide.shapes.add_textbox(Inches(1), Inches(y_pos), Inches(11.3), Inches(0.5))
                ttf = tb.text_frame
                tp = ttf.paragraphs[0]
                tp.text = line[:200]
                tp.font.size = Pt(16)
                tp.font.color.rgb = RGBColor(*has_text_color)
                tp.font.name = "Calibri"
                y_pos += 0.55

        num_box = slide.shapes.add_textbox(Inches(12), Inches(7), Inches(1), Inches(0.4))
        ntf = num_box.text_frame
        np = ntf.paragraphs[0]
        np.text = str(slide_idx + 2)
        np.font.size = Pt(10)
        np.font.color.rgb = RGBColor(*info_color["text_secondary"])
        np.alignment = PP_ALIGN.RIGHT

    # ── CLOSING SLIDE ──
    closing_keyword = "thank you elegant background"
    closing_image = buscar_y_descargar_imagen(closing_keyword)
    closing_info = analizar_paleta_y_brillo(closing_image) if closing_image else {
        "accent": _hex_to_rgb(pal["secondary"]),
        "text_main": _hex_to_rgb(pal["light"]),
        "text_secondary": (180, 180, 180),
        "is_background_light": False,
    }

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    if closing_image:
        slide.shapes.add_picture(closing_image, 0, 0, prs.slide_width, prs.slide_height)

    if not closing_image or closing_info["is_background_light"]:
        _overlay(slide)

    thanks_box = slide.shapes.add_textbox(Inches(1.5), Inches(2.5), Inches(10), Inches(2))
    tf = thanks_box.text_frame
    tp = tf.paragraphs[0]
    tp.text = "Gracias"
    tp.font.size = Pt(52)
    tp.font.bold = True
    tp.font.color.rgb = RGBColor(*closing_info["text_main"])
    tp.font.name = "Calibri Light"
    tp.alignment = PP_ALIGN.CENTER

    agata_box = slide.shapes.add_textbox(Inches(1.5), Inches(4.5), Inches(10), Inches(1))
    atf = agata_box.text_frame
    ap = atf.paragraphs[0]
    ap.text = "Disenado con amor por Agata  |  JARVIS Industries"
    ap.font.size = Pt(16)
    ap.font.color.rgb = RGBColor(*closing_info["accent"])
    ap.font.italic = True
    ap.alignment = PP_ALIGN.CENTER

    prs.save(str(filepath))
    return str(filepath)


def _hex_to_rgb(h: str) -> tuple:
    h = h.lstrip("#")
    return tuple(int(h[i : i + 2], 16) for i in (0, 2, 4))


def agata_create(parameters: dict, player=None, speak=None):
    doc_type = parameters.get("type", "word").lower()
    title = parameters.get("title", "Documento Agata")
    topic = parameters.get("topic", "")
    palette = parameters.get("palette", "elegant")
    api_key = parameters.get("api_key", "")

    if doc_type not in ("word", "ppt"):
        doc_type = "word"

    if not topic:
        return "Necesito un tema (topic) para crear el documento. Por favor, dime de que quieres que trate."

    player.write_log(f"[Agata] Creando {doc_type.upper()}: {title}")

    prompt = (
        f"Crea contenido profesional y bien estructurado para un {'documento Word' if doc_type == 'word' else 'presentacion PowerPoint'} "
        f"titulado '{title}' sobre el tema: {topic}.\n\n"
    )
    if doc_type == "word":
        prompt += (
            "Formato del contenido (usa estos marcadores para elementos especiales):\n\n"
            "# Titulo Principal\n"
            "Introduccion elegante y profesional (2-3 parrafos)\n\n"
            "## Seccion Principal\n"
            "Contenido detallado con datos relevantes y analisis profundo\n"
            "- Punto clave 1 con detalle\n"
            "- Punto clave 2 con detalle\n\n"
            "Para insertar una imagen usa: [IMAGE: keyword en ingles]\n"
            "Ejemplo: [IMAGE: data analysis dashboard]\n\n"
            "Para insertar una tabla usa:\n"
            "[TABLE]\n"
            "| Columna 1 | Columna 2 | Columna 3 |\n"
            "| Dato 1 | Dato 2 | Dato 3 |\n"
            "| Dato 4 | Dato 5 | Dato 6 |\n"
            "[ENDTABLE]\n\n"
            "Para citas destacadas usa: > Texto de la cita\n"
            "Para separadores de seccion usa: ---\n\n"
            "### Sub-seccion\n"
            "Detalles complementarios\n\n"
            "## Conclusion\n"
            "Resumen elegante y cierre profesional\n\n"
            "Incluye al menos 1 tabla con datos relevantes y 1 imagen si el tema lo permite.\n"
            "Escribe todo en espanol, lenguaje profesional pero calido.\n"
            "IMPORTANTE: No incluyas tu proceso de pensamiento ni analisis. "
            "Empieza directamente con el contenido del documento."
        )
    else:
        prompt += (
            "Formato: crea contenido por diapositivas separadas por '---'.\n"
            "Para cada diapositiva:\n"
            "## Titulo de Diapositiva [keyword: terminos de busqueda en ingles]\n"
            "**Idea Principal**\n"
            "- Punto clave 1\n"
            "- Punto clave 2\n"
            "- Punto clave 3\n"
            "Texto complementario breve\n\n"
            "IMPORTANTE: Despues de cada titulo de diapositiva, incluye SIEMPRE [keyword: ...] "
            "con 2-4 palabras clave en INGLES que describan una imagen de fondo ideal para esa diapositiva. "
            "Ejemplo: ## Exploracion Espacial [keyword: mars rover space exploration]\n"
            "Incluye 5-8 diapositivas. Escribe todo en espanol, lenguaje profesional pero calido.\n"
            "IMPORTANTE: No incluyas tu proceso de pensamiento ni analisis. "
            "Empieza directamente con el contenido de la primera diapositiva."
        )

    player.write_log("[Agata] Generando contenido con DeepSeek...")
    content = _call_opencode(prompt, max_tokens=4000)

    if content.startswith("[OpenCode") or content.startswith("[DeepSeek"):
        player.write_log(f"[Agata] OpenCode fallo: {content[:120]}. Usando Ollama (phi:2.7b)...")
        if speak:
            pass
        content = _call_ollama(prompt, max_tokens=4000)

    if content.startswith("[Ollama"):
        player.write_log(f"[Agata] Ollama fallo: {content[:120]}. Usando Gemini...")
        content = _call_gemini_fast(prompt, api_key)
        if content.startswith("[Gemini"):
            return f"No pude generar el contenido. Error: {content}"

    player.write_log(f"[Agata] Disenando {doc_type.upper()} con paleta {palette}...")

    try:
        if doc_type == "word":
            filepath = _design_word(content, title, palette, player, speak)
        else:
            filepath = _design_ppt(content, title, palette, player, speak)
    except Exception as e:
        return f"Error al crear el archivo: {str(e)[:200]}"

    if os.path.exists(filepath):
        size_kb = os.path.getsize(filepath) / 1024
        result = (
            f"Tu {'documento Word' if doc_type == 'word' else 'presentacion PowerPoint'} "
            f"'{title}' esta listo!\n"
            f"Archivo: {Path(filepath).name}\n"
            f"Tamaño: {size_kb:.1f} KB\n"
            f"Paleta: {palette}\n"
            f"Ubicacion: OneDrive > Documenti"
        )
        player.write_log(f"[Agata] Listo: {Path(filepath).name} ({size_kb:.1f} KB)")

        try:
            os.startfile(filepath)
        except Exception:
            pass

        if speak:
            pass

        return result
    else:
        return f"El archivo se creo pero no pude encontrarlo en: {filepath}"


def list_palettes(parameters=None, player=None, speak=None):
    result = "Paletas de diseno disponibles:\n\n"
    for name, pal in PALETTES.items():
        result += f"[{name.upper()}]\n"
        result += f"  Primario:  {pal['primary']}\n"
        result += f"  Secundario: {pal['secondary']}\n"
        result += f"  Acento:    {pal['accent']}\n\n"
    return result
